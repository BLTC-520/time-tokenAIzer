#!/usr/bin/env python3
"""Generate the TIME v4 hook explanation PPTX + PDF render."""
import os
import sys
from pathlib import Path

# Generation-only dependencies are intentionally kept out of package.json.
# Bootstrap when needed with:
#   python3 -m pip install --target /tmp/time-v4-artifact-py -r docs/presentations/requirements-time-v4-hook.txt
DEPS_DIRS = [
    Path(os.environ.get("TIME_V4_ARTIFACT_PY_DEPS", "")),
    Path(__file__).with_name(".time-v4-artifact-py"),
    Path("/tmp/time-v4-artifact-py"),
]
for deps_dir in DEPS_DIRS:
    if str(deps_dir) and deps_dir.exists():
        sys.path.insert(0, str(deps_dir))

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import landscape
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
except ModuleNotFoundError as exc:
    raise SystemExit(
        "Missing python-pptx/reportlab dependencies. Run: "
        "python3 -m pip install --target /tmp/time-v4-artifact-py "
        "-r docs/presentations/requirements-time-v4-hook.txt"
    ) from exc

OUT_DIR = Path("docs/presentations")
OUT_DIR.mkdir(parents=True, exist_ok=True)
PPTX_OUT = OUT_DIR / "time-v4-hook-implementation.pptx"
PDF_OUT = OUT_DIR / "time-v4-hook-implementation.pdf"

W_IN, H_IN = 13.333, 7.5
W_PT, H_PT = W_IN * inch, H_IN * inch

COLORS = {
    "navy": "0F172A",
    "slate": "334155",
    "muted": "64748B",
    "ice": "E0F2FE",
    "blue": "2563EB",
    "green": "16A34A",
    "green_bg": "DCFCE7",
    "purple": "7C3AED",
    "purple_bg": "EDE9FE",
    "amber": "D97706",
    "amber_bg": "FEF3C7",
    "red": "DC2626",
    "red_bg": "FEE2E2",
    "cream": "FFF7ED",
    "white": "FFFFFF",
    "off": "F8FAFC",
}

def rgb(hexstr):
    hexstr = hexstr.replace("#", "")
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))

def hx(name):
    return COLORS[name]

def add_bg(slide, color="off"):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(hx(color))

def ppt_text(slide, text, x, y, w, h, size=18, color="navy", bold=False, align="left", font="Aptos", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(hx(color) if color in COLORS else color)
    return box

def ppt_card(slide, x, y, w, h, title, body="", fill="white", line="slate", title_color="navy", body_color="slate"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(hx(fill) if fill in COLORS else fill)
    shape.line.color.rgb = rgb(hx(line) if line in COLORS else line)
    shape.line.width = Pt(1.3)
    ppt_text(slide, title, x+0.16, y+0.12, w-0.32, 0.34, 15.5, title_color, True)
    if body:
        ppt_text(slide, body, x+0.16, y+0.52, w-0.32, h-0.62, 12.3, body_color)
    return shape

def ppt_line(slide, x1, y1, x2, y2, color="slate", width=2.0):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(hx(color) if color in COLORS else color)
    line.line.width = Pt(width)
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass
    return line

def title(slide, heading, sub=None, dark=False):
    if dark:
        add_bg(slide, "navy")
        ppt_text(slide, heading, 0.65, 0.55, 9.7, 0.65, 32, "white", True)
        if sub:
            ppt_text(slide, sub, 0.68, 1.25, 10.8, 0.55, 17, "ice")
    else:
        add_bg(slide, "off")
        ppt_text(slide, heading, 0.6, 0.38, 10.8, 0.55, 28, "navy", True)
        if sub:
            ppt_text(slide, sub, 0.64, 0.98, 11.0, 0.4, 14, "slate")

def add_footer(slide, n):
    ppt_text(slide, f"time-tokenAIzer · Uniswap v4 hook artifacts · {n}/8", 0.6, 7.05, 6.2, 0.25, 9, "muted")

prs = Presentation()
prs.slide_width = Inches(W_IN)
prs.slide_height = Inches(H_IN)
blank = prs.slide_layouts[6]

# Slide 1
s = prs.slides.add_slide(blank)
title(s, "TIME v4 Hook Implementation", "Booking intent guard + telemetry for a time-credit marketplace", dark=True)
ppt_text(s, "Core thesis", 0.75, 2.15, 1.7, 0.3, 14, "ice", True)
ppt_text(s, "Uniswap v4 validates and observes booking intent during swap. BookingManager owns service rights, inventory, quote replay protection, settlement, and lifecycle.", 0.75, 2.55, 8.2, 1.15, 25, "white", True)
ppt_card(s, 9.3, 2.05, 3.25, 1.0, "Hook role", "beforeSwap guard\nafterSwap telemetry\nZERO_DELTA", "purple_bg", "purple", "purple")
ppt_card(s, 9.3, 3.28, 3.25, 1.0, "Booking role", "inventory + slots\nquote validity\nburn/lock TIME", "red_bg", "red", "red")
ppt_card(s, 9.3, 4.52, 3.25, 1.0, "Frontend flow", "quote → hookData → swap\nthen bookWithCredits", "green_bg", "green", "green")
add_footer(s, 1)

# Slide 2
s = prs.slides.add_slide(blank)
title(s, "Current implementation at a glance", "Grounded in the current contracts, frontend services, and docs.")
ppt_card(s, 0.75, 1.65, 3.6, 3.85, "TimePoolHook.sol", "- Enables only beforeSwap + afterSwap\n- Validates pool/router allowlists\n- Decodes booking hookData\n- Delegates quote/inventory validity to BookingManager\n- Emits TimeSwapObserved after swap", "purple_bg", "purple", "purple")
ppt_card(s, 4.85, 1.65, 3.6, 3.85, "BookingManager.sol", "- Provider inventory + pause state\n- EIP-712 quote validation\n- Quote replay protection via usedQuotes\n- Slot lock via slotTaken\n- bookWithCredits burns TIME and creates booking", "red_bg", "red", "red")
ppt_card(s, 8.95, 1.65, 3.6, 3.85, "Frontend + v4 SDK", "- /api/booking/quote returns BookingQuote\n- buildHookData ABI-encodes quote\n- Universal Router executes V4_SWAP\n- UI calls bookWithCredits after swap confirmation", "green_bg", "green", "green")
ppt_text(s, "Primary evidence: contracts/src/TimePoolHook.sol, contracts/src/BookingManager.sol, src/app/services/uniswapV4Service.ts, src/app/components/time-market/BookingCheckout.tsx", 0.8, 6.0, 11.9, 0.45, 12.5, "muted")
add_footer(s, 2)

# Slide 3
s = prs.slides.add_slide(blank)
title(s, "Hook permission profile", "The hook is intentionally narrow: no custom accounting and no booking lifecycle mutation.")
ppt_card(s, 0.8, 1.7, 2.9, 1.15, "beforeSwap", "ENABLED\nGuard: pool, router, hookData, quote, inventory", "purple_bg", "purple", "purple")
ppt_card(s, 4.0, 1.7, 2.9, 1.15, "afterSwap", "ENABLED\nTelemetry only: TimeSwapObserved", "purple_bg", "purple", "purple")
ppt_card(s, 7.2, 1.7, 2.35, 1.15, "Return deltas", "DISABLED\nNo custom swap amounts", "red_bg", "red", "red")
ppt_card(s, 9.85, 1.7, 2.35, 1.15, "Liquidity callbacks", "DISABLED\nNo hook-owned LP lifecycle", "red_bg", "red", "red")
# matrix
headers = ["Callback family", "Status", "Why it matters"]
xs = [0.85, 3.65, 5.4]
for i,h in enumerate(headers): ppt_text(s,h,xs[i],3.35,[2.3,1.3,6.2][i],0.3,13,"navy",True)
rows = [
    ("beforeSwap / afterSwap", "On", "Validate booking-aware intent during swap and emit lightweight observability."),
    ("beforeSwapReturnDelta / afterSwapReturnDelta", "Off", "Avoids NoOp/custom-delta rug and prevents hook from becoming AMM logic."),
    ("add/remove liquidity callbacks", "Off", "Keeps LP accounting outside the marketplace hook."),
    ("booking settlement", "Not a hook permission", "Settlement lives in BookingManager after final TIME ownership exists."),
]
y=3.85
for idx,row in enumerate(rows):
    fill = "white" if idx%2==0 else "ice"
    ppt_card(s,0.75,y,11.8,0.55,"", "", fill, "ice")
    ppt_text(s,row[0],0.9,y+0.12,2.55,0.25,11.5,"slate")
    ppt_text(s,row[1],3.7,y+0.12,1.1,0.25,11.5,"green" if row[1]=='On' else 'red', True)
    ppt_text(s,row[2],5.1,y+0.12,6.9,0.3,11.2,"slate")
    y += 0.65
add_footer(s, 3)

# Slide 4 activity
s = prs.slides.add_slide(blank)
title(s, "Activity flow: swap then book", "A valid booking-aware swap does not automatically create a booking.")
steps = [
    ("1", "Select provider + slot", "Buyer chooses service inventory in the UI", "blue"),
    ("2", "Get signed quote", "API/BookingManager checks provider, hours, slot", "green"),
    ("3", "Encode hookData", "Frontend builds ABI quote tuple", "green"),
    ("4", "V4 swap executes", "Universal Router → PoolManager → TIME/USDC pool", "amber"),
    ("5", "Hook validates intent", "beforeSwap checks allowlists, quote, inventory, replay policy", "purple"),
    ("6", "Telemetry emitted", "afterSwap emits TimeSwapObserved only", "purple"),
    ("7", "Book with credits", "BookingManager burns TIME, locks slot, creates booking", "red"),
]
for i,(num,head,body,color) in enumerate(steps):
    x = 0.7 + (i%4)*3.05; y = 1.65 + (i//4)*2.05
    ppt_card(s,x,y,2.62,1.25,f"{num}. {head}",body, f"{color}_bg" if color in ['green','purple','red','amber'] else 'ice', color, color)
    if i < 6:
        if i%4 != 3:
            ppt_line(s, x+2.62, y+0.64, x+2.98, y+0.64, "muted", 1.4)
        else:
            ppt_line(s, x+1.31, y+1.25, 0.7, y+2.05, "muted", 1.4)
ppt_card(s, 9.65, 5.58, 2.95, 1.0, "Failure mode", "If swap succeeds but booking fails, UI can retry bookWithCredits with settled TIME.", "cream", "amber", "amber")
add_footer(s, 4)

# Slide 5 architecture
s = prs.slides.add_slide(blank)
title(s, "System architecture", "Marketplace/app domain stays separate from Uniswap v4 liquidity domain.")
# domain backgrounds
ppt_card(s, 0.55, 1.38, 5.65, 4.95, "Marketplace / app domain", "", "ice", "blue", "blue")
ppt_card(s, 6.75, 1.38, 5.95, 4.95, "Uniswap v4 liquidity domain", "", "purple_bg", "purple", "purple")
# components
ppt_card(s,0.9,2.1,1.65,0.78,"Buyer wallet","USDC + TIME","white","blue","blue")
ppt_card(s,3.15,1.95,2.5,0.95,"Next.js frontend","BookingCheckout + marketplace UI","white","blue","blue")
ppt_card(s,3.15,3.22,2.5,0.95,"/api/booking/quote","real EIP-712 quote\nmock UI-only warning","green_bg","green","green")
ppt_card(s,3.05,4.7,2.65,1.0,"BookingManager","inventory, slots, usedQuotes, bookWithCredits","red_bg","red","red")
ppt_card(s,0.95,4.9,1.65,0.8,"TIME token","1e18 = 1 hour","amber_bg","amber","amber")
ppt_card(s,7.15,2.05,2.05,0.86,"Universal Router","Permit2 + V4_SWAP","white","purple","purple")
ppt_card(s,10.0,2.05,2.05,0.86,"PoolManager","flash accounting callbacks","white","purple","purple")
ppt_card(s,10.0,3.65,2.05,0.86,"TIME/USDC pool","price discovery","amber_bg","amber","amber")
ppt_card(s,7.1,3.75,2.28,1.12,"TimePoolHook","beforeSwap guard\nafterSwap telemetry\nZERO_DELTA","purple_bg","purple","purple")
ppt_card(s,7.25,5.48,1.95,0.62,"Event stream","TimeSwapObserved","white","muted","slate")
for a in [(2.55,2.48,3.15,2.43),(4.4,2.9,4.4,3.22),(4.4,4.17,4.4,4.7),(5.65,2.43,7.15,2.45),(9.2,2.48,10.0,2.48),(11.0,2.91,11.0,3.65),(10.0,4.05,9.38,4.25),(7.1,4.35,5.7,5.15),(8.2,4.87,8.2,5.48),(7.15,2.9,5.7,4.95)]:
    ppt_line(s,*a,"muted",1.3)
ppt_text(s,"Invariant: booking state is never created by the hook; it is created only in BookingManager after settled TIME credits are available.",0.9,6.45,11.4,0.38,12.3,"red",True)
add_footer(s, 5)

# Slide 6 security
s = prs.slides.add_slide(blank)
title(s, "Security boundary and invariants", "The presentation should not imply atomic service booking inside a v4 callback.")
items = [
    ("No booking in callbacks", "Hook never calls bookWithCredits, never burns TIME, never reserves a slot.", "red"),
    ("Pool + router allowlists", "_validatePoolAndRouter gates both beforeSwap and afterSwap.", "purple"),
    ("hookData policy is explicit", "Empty hookData can pass unless strict mode is enabled; mock quotes are UI/dev only.", "amber"),
    ("Quote replay is split", "Hook-level consumedHookQuote guards trusted routers; BookingManager.usedQuotes creates booking replay protection.", "green"),
    ("afterSwap is telemetry only", "It decodes buyer/quoteId and emits TimeSwapObserved; it returns zero settlement delta.", "purple"),
    ("Booking can fail after swap", "UI preserves a retry state when swap confirmation succeeds but booking transaction fails.", "red"),
]
for i,(h,b,c) in enumerate(items):
    x = 0.75 + (i%2)*6.05; y = 1.55 + (i//2)*1.55
    ppt_card(s,x,y,5.55,1.05,h,b,f"{c}_bg" if c in ['green','purple','red','amber'] else 'white',c,c)
add_footer(s, 6)

# Slide 7 verification
s = prs.slides.add_slide(blank)
title(s, "Verification evidence in the repo", "Existing tests and docs support the artifact claims.")
ppt_card(s,0.75,1.55,3.65,4.35,"Contract tests", "TimePoolHook.t.sol covers permissions, owner config, pool/router allowlists, malformed hookData, missing buyer, zero hours, inventory bounds, invalid quote, quote replay, untrusted routers, and telemetry-only afterSwap.", "purple_bg", "purple", "purple")
ppt_card(s,4.85,1.55,3.65,4.35,"Source files", "TimePoolHook.sol shows beforeSwap/afterSwap only and ZERO_DELTA. BookingManager.sol shows usedQuotes, inventory decrement, slotTaken, TIME burn, and Booking struct creation inside bookWithCredits.", "red_bg", "red", "red")
ppt_card(s,8.95,1.55,3.65,4.35,"Docs + app flow", "Architecture and security docs record the boundary. BookingCheckout runs swap first, then submitBooking; failure state supports retry after swap confirmation.", "green_bg", "green", "green")
ppt_text(s,"Useful commands for deeper verification: forge test -vvv; forge snapshot --match-contract TimePoolHookTest --skip script; npm run build; npx tsc --noEmit --pretty false --incremental false",0.8,6.25,11.8,0.4,11.8,"muted")
add_footer(s, 7)

# Slide 8
s = prs.slides.add_slide(blank)
title(s, "Simple explanation for demos", "Use this framing when walking non-specialists through the system.", dark=True)
ppt_card(s,0.8,1.75,3.65,2.3,"1. TIME is liquid", "Users can acquire fungible TIME credits through a Uniswap v4 TIME/USDC pool.", "white", "blue", "blue")
ppt_card(s,4.85,1.75,3.65,2.3,"2. Hook checks intent", "The hook can reject unsafe booking-aware swaps but does not create service rights.", "white", "purple", "purple")
ppt_card(s,8.9,1.75,3.65,2.3,"3. Booking is separate", "BookingManager burns/locks settled TIME and creates the actual booking record.", "white", "red", "red")
ppt_text(s,"Deliverables",0.9,4.85,1.6,0.3,15,"ice",True)
ppt_text(s,"- docs/diagrams/time-v4-booking-activity.excalidraw\n- docs/diagrams/time-v4-system-architecture.excalidraw\n- docs/presentations/time-v4-hook-implementation.pptx\n- docs/presentations/time-v4-hook-implementation.pdf",0.9,5.25,8.4,1.0,14,"white")
add_footer(s, 8)

prs.save(PPTX_OUT)

# PDF render helpers
try:
    pdfmetrics.registerFont(TTFont('AptosLike', '/System/Library/Fonts/Supplemental/Arial.ttf'))
except Exception:
    pass

FONT = 'AptosLike' if 'AptosLike' in pdfmetrics.getRegisteredFontNames() else 'Helvetica'
FONT_BOLD = 'Helvetica-Bold'

def ccol(name):
    return colors.HexColor('#' + (COLORS[name] if name in COLORS else name))

def wrap_lines(text, width_chars=42):
    out=[]
    for para in text.split('\n'):
        words=para.split()
        if not words:
            out.append('')
            continue
        line=''
        for w in words:
            if len(line)+len(w)+1 > width_chars:
                out.append(line)
                line=w
            else:
                line = w if not line else line+' '+w
        out.append(line)
    return out

def draw_text(c, text, x, y, w, h, size=14, color='navy', bold=False, align='left'):
    c.setFillColor(ccol(color) if color in COLORS else colors.HexColor('#'+color))
    c.setFont(FONT_BOLD if bold else FONT, size)
    lines = wrap_lines(text, max(12, int(w/(size*0.45))))
    line_h = size*1.2
    cy = H_PT - y - size
    for line in lines[:max(1,int(h/line_h))]:
        if align == 'center':
            c.drawCentredString(x+w/2, cy, line)
        else:
            c.drawString(x, cy, line)
        cy -= line_h

def draw_card(c, x, y, w, h, title, body='', fill='white', line='slate', title_color='navy', body_color='slate'):
    c.setFillColor(ccol(fill)); c.setStrokeColor(ccol(line)); c.setLineWidth(1.1)
    c.roundRect(x, H_PT-y-h, w, h, 10, fill=1, stroke=1)
    if title:
        draw_text(c, title, x+12, y+12, w-24, 24, 12.5, title_color, True)
    if body:
        draw_text(c, body, x+12, y+44, w-24, h-50, 9.8, body_color)

def draw_line(c, x1,y1,x2,y2,color='muted'):
    c.setStrokeColor(ccol(color)); c.setLineWidth(1.2)
    c.line(x1,H_PT-y1,x2,H_PT-y2)

pdf = canvas.Canvas(str(PDF_OUT), pagesize=(W_PT,H_PT))
# To keep PDF readable, mirror the exact slide narrative with simpler reportlab rendering.
for idx in range(1,9):
    pdf.setFillColor(ccol('navy' if idx in [1,8] else 'off'))
    pdf.rect(0,0,W_PT,H_PT,stroke=0,fill=1)
    dark = idx in [1,8]
    if idx == 1:
        draw_text(pdf, 'TIME v4 Hook Implementation', 45, 42, 620, 50, 27, 'white', True)
        draw_text(pdf, 'Booking intent guard + telemetry for a time-credit marketplace', 48, 96, 760, 32, 15, 'ice')
        draw_text(pdf, 'Uniswap v4 validates and observes booking intent during swap. BookingManager owns service rights, inventory, quote replay protection, settlement, and lifecycle.', 55, 185, 610, 95, 22, 'white', True)
        draw_card(pdf, 670, 150, 235, 95, 'Hook role', 'beforeSwap guard\nafterSwap telemetry\nZERO_DELTA', 'purple_bg', 'purple', 'purple')
        draw_card(pdf, 670, 265, 235, 95, 'Booking role', 'inventory + slots\nquote validity\nburn/lock TIME', 'red_bg', 'red', 'red')
        draw_card(pdf, 670, 380, 235, 82, 'Frontend flow', 'quote → hookData → swap\nthen bookWithCredits', 'green_bg', 'green', 'green')
    elif idx == 2:
        draw_text(pdf,'Current implementation at a glance',45,34,760,40,23,'navy',True)
        draw_text(pdf,'Grounded in the current contracts, frontend services, and docs.',48,80,770,26,12,'slate')
        draw_card(pdf,55,125,260,290,'TimePoolHook.sol','- Enables only beforeSwap + afterSwap\n- Validates pool/router allowlists\n- Decodes booking hookData\n- Delegates quote/inventory validity\n- Emits TimeSwapObserved', 'purple_bg','purple','purple')
        draw_card(pdf,350,125,260,290,'BookingManager.sol','- Provider inventory + pause\n- EIP-712 quote validation\n- Replay protection via usedQuotes\n- Slot lock via slotTaken\n- bookWithCredits burns TIME', 'red_bg','red','red')
        draw_card(pdf,645,125,260,290,'Frontend + v4 SDK','- Quote API returns BookingQuote\n- buildHookData ABI-encodes quote\n- Universal Router executes V4_SWAP\n- UI calls bookWithCredits after swap', 'green_bg','green','green')
        draw_text(pdf,'Evidence: TimePoolHook.sol, BookingManager.sol, uniswapV4Service.ts, BookingCheckout.tsx',60,450,850,26,10,'muted')
    elif idx == 3:
        draw_text(pdf,'Hook permission profile',45,34,760,40,23,'navy',True)
        draw_text(pdf,'The hook is intentionally narrow: no custom accounting and no booking lifecycle mutation.',48,80,800,26,12,'slate')
        for x,t,b,f,l in [(55,'beforeSwap','ENABLED\nGuard + validation','purple_bg','purple'),(285,'afterSwap','ENABLED\nTelemetry only','purple_bg','purple'),(515,'Return deltas','DISABLED\nNo custom swap amounts','red_bg','red'),(715,'Liquidity callbacks','DISABLED\nNo LP lifecycle','red_bg','red')]:
            draw_card(pdf,x,130,180,82,t,b,f,l,l)
        draw_card(pdf,55,260,850,55,'beforeSwap / afterSwap: On','Validate booking-aware intent and emit observability.','white','ice')
        draw_card(pdf,55,325,850,55,'Return deltas: Off','Avoids custom-delta/NoOp risk and keeps hook out of AMM accounting.','ice','ice')
        draw_card(pdf,55,390,850,45,'Liquidity callbacks: Off','Keeps LP accounting outside the marketplace hook.','white','ice')
        draw_card(pdf,55,445,850,45,'Booking settlement: Not a hook permission','Settlement lives in BookingManager after final TIME ownership exists.','ice','ice')
    elif idx == 4:
        draw_text(pdf,'Activity flow: swap then book',45,34,760,40,23,'navy',True)
        draw_text(pdf,'A valid booking-aware swap does not automatically create a booking.',48,80,770,26,12,'slate')
        cards=[('1. Select provider + slot','Buyer chooses service inventory','ice','blue'),('2. Get real quote','EIP-712 quote; mock is UI-only','green_bg','green'),('3. Encode hookData','Frontend ABI-encodes quote','green_bg','green'),('4. V4 swap executes','Universal Router -> PoolManager','amber_bg','amber'),('5. Hook validates intent','beforeSwap checks quote + inventory','purple_bg','purple'),('6. Telemetry emitted','afterSwap emits event only','purple_bg','purple'),('7. Book with credits','BookingManager burns TIME + locks slot','red_bg','red')]
        for i,(h,b,f,l) in enumerate(cards):
            x=55+(i%4)*220; y=130+(i//4)*155
            draw_card(pdf,x,y,190,90,h,b,f,l,l)
        draw_card(pdf,690,435,230,82,'Failure mode','Swap OK + booking fail => retry bookWithCredits with settled TIME.','cream','amber','amber')
    elif idx == 5:
        draw_text(pdf,'System architecture',45,34,760,40,23,'navy',True)
        draw_text(pdf,'Marketplace/app domain stays separate from Uniswap v4 liquidity domain.',48,80,770,26,12,'slate')
        draw_card(pdf,45,120,405,330,'Marketplace / app domain','', 'ice','blue','blue')
        draw_card(pdf,500,120,420,330,'Uniswap v4 liquidity domain','', 'purple_bg','purple','purple')
        draw_card(pdf,75,185,120,55,'Buyer wallet','USDC + TIME','white','blue','blue')
        draw_card(pdf,235,170,170,70,'Frontend','BookingCheckout UI','white','blue','blue')
        draw_card(pdf,235,270,170,78,'Quote API','Real EIP-712 quote\nmock UI-only warning','green_bg','green','green')
        draw_card(pdf,230,380,185,70,'BookingManager','inventory, slots, bookWithCredits','red_bg','red','red')
        draw_card(pdf,535,185,150,60,'Universal Router','Permit2 + V4_SWAP','white','purple','purple')
        draw_card(pdf,740,185,150,60,'PoolManager','callbacks','white','purple','purple')
        draw_card(pdf,740,305,150,60,'TIME/USDC pool','price discovery','amber_bg','amber','amber')
        draw_card(pdf,535,310,165,80,'TimePoolHook','beforeSwap guard\nafterSwap telemetry','purple_bg','purple','purple')
        draw_text(pdf,'Invariant: booking state is created only by BookingManager after settled TIME credits are available.',65,470,800,22,10,'red',True)
    elif idx == 6:
        draw_text(pdf,'Security boundary and invariants',45,34,760,40,23,'navy',True)
        draw_text(pdf,'Do not imply atomic service booking inside a v4 callback.',48,80,770,26,12,'slate')
        items=[('No booking in callbacks','Hook never calls bookWithCredits, burns TIME, or reserves slots.','red'),('Pool + router allowlists','Both callbacks call _validatePoolAndRouter.','purple'),('hookData policy','Empty hookData can pass unless strict; mock is UI/dev only.','amber'),('Quote replay is split','Hook consumedHookQuote is not BookingManager.usedQuotes.','green'),('afterSwap telemetry only','It emits TimeSwapObserved and returns zero delta.','purple'),('Booking can fail after swap','UI supports retry with settled TIME.','red')]
        for i,(h,b,c) in enumerate(items):
            x=55+(i%2)*430; y=125+(i//2)*105
            draw_card(pdf,x,y,375,70,h,b,c+'_bg',c,c)
    elif idx == 7:
        draw_text(pdf,'Verification evidence in the repo',45,34,760,40,23,'navy',True)
        draw_text(pdf,'Existing tests and docs support the artifact claims.',48,80,770,26,12,'slate')
        draw_card(pdf,55,125,260,290,'Contract tests','TimePoolHook.t.sol covers permissions, owner config, allowlists, malformed hookData, buyer/hours validation, inventory bounds, invalid quote, quote replay, untrusted routers, and afterSwap telemetry.','purple_bg','purple','purple')
        draw_card(pdf,350,125,260,290,'Source files','TimePoolHook.sol shows beforeSwap/afterSwap only and ZERO_DELTA. BookingManager.sol shows usedQuotes, inventory decrement, slotTaken, TIME burn, and booking creation.','red_bg','red','red')
        draw_card(pdf,645,125,260,290,'Docs + app flow','Architecture and security docs record the boundary. BookingCheckout runs swap first, then submitBooking.','green_bg','green','green')
    elif idx == 8:
        draw_text(pdf,'Simple explanation for demos',45,42,720,42,25,'white',True)
        draw_text(pdf,'Use this framing when walking non-specialists through the system.',48,90,770,26,13,'ice')
        draw_card(pdf,55,150,260,150,'1. TIME is liquid','Users can acquire fungible TIME credits through a Uniswap v4 TIME/USDC pool.','white','blue','blue')
        draw_card(pdf,350,150,260,150,'2. Hook checks intent','The hook can reject unsafe booking-aware swaps but does not create service rights.','white','purple','purple')
        draw_card(pdf,645,150,260,150,'3. Booking is separate','BookingManager burns/locks settled TIME and creates the actual booking record.','white','red','red')
        draw_text(pdf,'Deliverables:\n- docs/diagrams/time-v4-booking-activity.excalidraw\n- docs/diagrams/time-v4-system-architecture.excalidraw\n- docs/presentations/time-v4-hook-implementation.pptx\n- docs/presentations/time-v4-hook-implementation.pdf',65,365,780,95,12,'white')
    draw_text(pdf, f'time-tokenAIzer · Uniswap v4 hook artifacts · {idx}/8', 45, 514, 420, 18, 7.5, 'ice' if dark else 'muted')
    pdf.showPage()
pdf.save()

print(PPTX_OUT, PPTX_OUT.stat().st_size)
print(PDF_OUT, PDF_OUT.stat().st_size)
